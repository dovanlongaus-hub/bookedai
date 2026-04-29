#!/usr/bin/env python3
"""
BookedAI 5-minute pitch deck builder.

10 slides, ~30 seconds each — visualization-heavy:
  - Solution slide drawn as a real flow diagram (channels → agent → ledger).
  - Live Proof embeds the production homepage screenshot + a phone-style
    Telegram conversation mockup rendered in shapes.
  - Tenant Cases embed the actual chess + AI Mentor hero artwork.
  - Why-We-Win uses a 2D quadrant scatter to put BookedAI alone top-right.

Run:
    cd pitch-deck
    ../.venv-pdf/bin/python build_5min.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

import content_5min as C
import theme as T


HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUTPUT = HERE / "output"


# ============================================================
# Low-level shape helpers (mirrors build_pitch_deck.py style)
# ============================================================

def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_no_fill(shape):
    shape.fill.background()


def _set_line(shape, color, width=Pt(0.75)):
    line = shape.line
    if color is None:
        line.fill.background()
        return
    line.color.rgb = color
    line.width = width


def add_rect(slide, l, t, w, h, *, fill=None, line=None, line_w=Pt(0.75),
             shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.shadow.inherit = False
    _set_fill(s, fill) if fill is not None else _set_no_fill(s)
    _set_line(s, line, line_w)
    return s


def add_round_rect(slide, l, t, w, h, *, fill=None, line=None, line_w=Pt(0.75)):
    return add_rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def add_oval(slide, l, t, w, h, *, fill=None, line=None, line_w=Pt(0.75)):
    return add_rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w,
                    shape=MSO_SHAPE.OVAL)


def add_text(slide, l, t, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, fill=None, line=None, margin=Inches(0)):
    box = slide.shapes.add_textbox(l, t, w, h)
    if fill is not None:
        _set_fill(box, fill)
    if line is not None:
        _set_line(box, line)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    tf.vertical_anchor = anchor

    paragraphs = runs if (runs and isinstance(runs[0], list)) else [runs]
    for p_index, paragraph in enumerate(paragraphs):
        para = tf.paragraphs[0] if p_index == 0 else tf.add_paragraph()
        para.alignment = align
        for r_index, r in enumerate(paragraph):
            run = (para.runs[0] if (r_index == 0 and p_index == 0 and para.runs)
                   else para.add_run())
            run.text = r.get("text", "")
            T.apply_text_style(
                run,
                size=r.get("size", T.SZ_BODY),
                color=r.get("color", T.INK),
                bold=r.get("bold", False),
                font=r.get("font", T.FONT_DISPLAY),
                tracking=r.get("tracking"),
            )
        last = paragraph[-1] if paragraph else {}
        if "space_after" in last:
            para.space_after = last["space_after"]
    return box


def fill_slide_bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, T.SLIDE_WIDTH, T.SLIDE_HEIGHT)
    _set_fill(s, color)
    _set_line(s, None)
    s.shadow.inherit = False
    spTree = s._element.getparent()
    spTree.remove(s._element)
    spTree.insert(2, s._element)
    return s


def add_image_safe(slide, path, l, t, *, width=None, height=None):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), l, t, width=width, height=height)


def add_kicker(slide, l, t, text, color=T.ACCENT, *, w=Inches(8)):
    return add_text(slide, l, t, w, Inches(0.32),
                    [{"text": text.upper(), "size": T.SZ_KICKER, "color": color,
                      "bold": True, "font": T.FONT_TEXT, "tracking": 240}])


def add_title(slide, l, t, text, *, w=Inches(11.5), color=T.INK,
              size=T.SZ_TITLE, bold=True):
    return add_text(slide, l, t, w, Inches(1.3),
                    [{"text": text, "size": size, "color": color, "bold": bold,
                      "font": T.FONT_DISPLAY, "tracking": -20}])


def add_brand_mark(slide, *, dark=False, l=None, t=Inches(0.45), h=Inches(0.42)):
    if l is None:
        l = T.MARGIN_X
    logo = ASSETS / ("logo-light.png" if dark else "logo-dark.png")
    if logo.exists():
        add_image_safe(slide, logo, l, t, height=h)


def add_slide_number(slide, n, total, *, dark=False):
    color = T.INK_INV_3 if dark else T.INK_3
    add_text(slide, T.SLIDE_WIDTH - Inches(1.4), T.SLIDE_HEIGHT - Inches(0.32),
             Inches(1.1), Inches(0.22),
             [{"text": f"{n:02d} / {total:02d}", "size": T.SZ_TINY,
               "color": color, "font": T.FONT_TEXT}], align=PP_ALIGN.RIGHT)


def _content_header(s, kicker, title, *, dark=False, kcolor=None):
    fill_slide_bg(s, T.BG_DARK if dark else T.BG_LIGHT)
    add_brand_mark(s, dark=dark)
    add_kicker(s, T.MARGIN_X, Inches(1.35), kicker,
               color=(kcolor if kcolor else (T.ACCENT_SOFT if dark else T.ACCENT)))
    add_title(s, T.MARGIN_X, Inches(1.7), title,
              color=T.INK_INV if dark else T.INK)


# ============================================================
# Slide 01 — Cover
# ============================================================

def slide_cover(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, T.BG_DARK)

    # Decorative gradient mark on the right
    mark = ASSETS / "mark-gradient.png"
    if mark.exists():
        add_image_safe(s, mark, T.SLIDE_WIDTH - Inches(5.5), Inches(0.5),
                       height=Inches(6.5))

    # Brand mark, top-left
    add_brand_mark(s, dark=True, t=Inches(0.55), h=Inches(0.5))

    # Hero — bold tagline
    add_text(s, T.MARGIN_X, Inches(2.0), Inches(9), Inches(1.1),
             [{"text": C.COMPANY, "size": Pt(72), "color": T.INK_INV,
               "bold": True, "tracking": -50}])
    add_text(s, T.MARGIN_X, Inches(3.3), Inches(11), Inches(0.85),
             [{"text": C.TAGLINE, "size": Pt(38), "color": T.ACCENT_SOFT,
               "bold": True, "tracking": -20}])
    add_text(s, T.MARGIN_X, Inches(4.3), Inches(11), Inches(0.55),
             [{"text": C.SUBTAGLINE, "size": T.SZ_SUBTITLE,
               "color": T.INK_INV_2, "font": T.FONT_TEXT}])

    # Hairline accent
    add_rect(s, T.MARGIN_X, Inches(5.4), Inches(1.5), Pt(2),
             fill=T.ACCENT, line=None)

    # Demo URL chip — eye-catching
    chip_w = Inches(5)
    chip_h = Inches(0.55)
    chip_y = Inches(5.7)
    add_round_rect(s, T.MARGIN_X, chip_y, chip_w, chip_h,
                   fill=T.ACCENT, line=None)
    add_text(s, T.MARGIN_X, chip_y, chip_w, chip_h,
             [
                 {"text": "→  Live demo: ", "size": T.SZ_BODY_S,
                  "color": T.INK_INV, "bold": True},
                 {"text": C.DEMO_URL, "size": T.SZ_BODY_S,
                  "color": T.INK_INV, "bold": True, "font": T.FONT_MONO},
             ], anchor=MSO_ANCHOR.MIDDLE, margin=Inches(0.3))

    # Audience + presenter footer block
    add_text(s, T.MARGIN_X, Inches(6.4), Inches(8), Inches(0.32),
             [{"text": C.AUDIENCE, "size": T.SZ_CAPTION,
               "color": T.INK_INV_2, "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X, Inches(6.7), Inches(8), Inches(0.32),
             [{"text": C.PRESENTER, "size": T.SZ_CAPTION,
               "color": T.INK_INV_3, "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X, Inches(7.0), Inches(8), Inches(0.3),
             [{"text": C.DECK_DATE.isoformat(), "size": T.SZ_TINY,
               "color": T.INK_INV_3, "font": T.FONT_MONO}])
    add_slide_number(s, n, total, dark=True)


# ============================================================
# Slide 02 — Problem (big stat hook)
# ============================================================

def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.PROBLEM["kicker"],
                    "Service SMEs lose more revenue to slow replies than to bad SEO.",
                    kcolor=T.RED)

    # Massive stat — left half
    stat_x = T.MARGIN_X
    stat_y = Inches(3.1)
    add_text(s, stat_x, stat_y, Inches(6), Inches(2.5),
             [{"text": C.PROBLEM["big_stat"], "size": Pt(180),
               "color": T.RED, "bold": True, "tracking": -100}])
    add_text(s, stat_x, stat_y + Inches(2.6), Inches(6.2), Inches(1.6),
             [{"text": C.PROBLEM["big_stat_label"],
               "size": T.SZ_SUBTITLE, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Pain bullets — right half
    pain_x = T.MARGIN_X + Inches(6.5)
    pain_y = Inches(3.0)
    pain_w = Inches(6.2)
    add_text(s, pain_x, pain_y, pain_w, Inches(0.32),
             [{"text": "WHERE THE REVENUE LEAKS", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    cur_y = pain_y + Inches(0.5)
    for i, (head, body) in enumerate(C.PROBLEM["pains"]):
        add_round_rect(s, pain_x, cur_y, pain_w, Inches(1.05),
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, pain_x, cur_y, Inches(0.06), Inches(1.05),
                 fill=T.RED, line=None)
        add_text(s, pain_x + Inches(0.28), cur_y + Inches(0.14),
                 pain_w - Inches(0.4), Inches(0.36),
                 [{"text": head, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, pain_x + Inches(0.28), cur_y + Inches(0.5),
                 pain_w - Inches(0.4), Inches(0.5),
                 [{"text": body, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])
        cur_y += Inches(1.2)

    add_slide_number(s, n, total)


# ============================================================
# Slide 03 — Solution (architecture flow diagram)
# ============================================================

def slide_solution(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.SOLUTION["kicker"], C.SOLUTION["title"])

    # Subtitle
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(12.13), Inches(0.4),
             [{"text": C.SOLUTION["subtitle"], "size": T.SZ_BODY,
               "color": T.INK_2, "font": T.FONT_TEXT}])

    # ----- Diagram -----
    # Layout: 3 horizontal bands
    #   Top (Inches 3.6-4.2): channel chips
    #   Middle (Inches 4.5-5.4): agent core (one wide rounded box w/ 3 sub-cells)
    #   Bottom (Inches 5.7-6.6): output ledger panels

    # ---- Channels row ----
    ch_y = Inches(3.55)
    ch_h = Inches(0.6)
    n_ch = len(C.SOLUTION["channels"])
    total_ch_w = Inches(12.13)
    gap = Inches(0.12)
    ch_w = (total_ch_w - gap * (n_ch - 1)) / n_ch
    for i, ch in enumerate(C.SOLUTION["channels"]):
        x = T.MARGIN_X + i * (ch_w + gap)
        add_round_rect(s, x, ch_y, ch_w, ch_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x, ch_y, ch_w, ch_h,
                 [{"text": ch, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- Down arrow ----
    arrow_y = ch_y + ch_h + Inches(0.05)
    add_text(s, T.SLIDE_WIDTH / 2 - Inches(0.5), arrow_y, Inches(1), Inches(0.3),
             [{"text": "▼", "size": Pt(18), "color": T.ACCENT, "bold": True}],
             align=PP_ALIGN.CENTER)

    # ---- Agent core (dark wide bar with 3 cells) ----
    agent_y = arrow_y + Inches(0.35)
    agent_h = Inches(1.05)
    agent_w = Inches(12.13)
    add_round_rect(s, T.MARGIN_X, agent_y, agent_w, agent_h,
                   fill=T.BG_DARK, line=T.ACCENT_DEEP, line_w=Pt(1))
    # Header label
    add_text(s, T.MARGIN_X + Inches(0.3), agent_y + Inches(0.08),
             Inches(6), Inches(0.3),
             [{"text": "BOOKEDAI AI AGENT LAYER  ·  shared messaging_automation_service",
               "size": T.SZ_TINY, "color": T.ACCENT_SOFT, "bold": True,
               "tracking": 220, "font": T.FONT_MONO}])

    # 3 sub-cells inside
    sub_n = len(C.SOLUTION["agents"])
    sub_y = agent_y + Inches(0.42)
    sub_h = Inches(0.55)
    sub_gap = Inches(0.18)
    sub_w = (agent_w - Inches(0.6) - sub_gap * (sub_n - 1)) / sub_n
    for i, (title, desc) in enumerate(C.SOLUTION["agents"]):
        x = T.MARGIN_X + Inches(0.3) + i * (sub_w + sub_gap)
        add_round_rect(s, x, sub_y, sub_w, sub_h,
                       fill=T.BG_PANEL_DARK, line=None)
        add_text(s, x + Inches(0.18), sub_y, sub_w - Inches(0.3), sub_h,
                 [
                     {"text": title + "  ", "size": T.SZ_BODY_S,
                      "color": T.INK_INV, "bold": True},
                     {"text": "›  " + desc, "size": T.SZ_CAPTION,
                      "color": T.INK_INV_2, "font": T.FONT_TEXT},
                 ], anchor=MSO_ANCHOR.MIDDLE)

    # ---- Down arrow ----
    arrow2_y = agent_y + agent_h + Inches(0.05)
    add_text(s, T.SLIDE_WIDTH / 2 - Inches(0.5), arrow2_y,
             Inches(1), Inches(0.3),
             [{"text": "▼", "size": Pt(18), "color": T.ACCENT, "bold": True}],
             align=PP_ALIGN.CENTER)

    # ---- Output panels ----
    out_y = arrow2_y + Inches(0.35)
    out_h = Inches(0.85)
    out_n = len(C.SOLUTION["outputs"])
    out_w = (agent_w - gap * (out_n - 1)) / out_n
    for i, (head, sub) in enumerate(C.SOLUTION["outputs"]):
        x = T.MARGIN_X + i * (out_w + gap)
        add_round_rect(s, x, out_y, out_w, out_h,
                       fill=T.BG_PANEL, line=T.ACCENT, line_w=Pt(1.5))
        add_text(s, x + Inches(0.25), out_y + Inches(0.16),
                 out_w - Inches(0.4), Inches(0.32),
                 [{"text": head, "size": T.SZ_BODY, "color": T.ACCENT, "bold": True}])
        add_text(s, x + Inches(0.25), out_y + Inches(0.5),
                 out_w - Inches(0.4), Inches(0.4),
                 [{"text": sub, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


# ============================================================
# Slide 04 — Live Proof (real screenshot + phone mockup)
# ============================================================

def slide_live_proof(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.LIVE_PROOF["kicker"], C.LIVE_PROOF["title"],
                    kcolor=T.GREEN)

    # ----- Left: real product screenshot in browser frame -----
    sshot = ASSETS / "product-after-search.png"  # 1280x720, perfect 16:9
    if not sshot.exists():
        sshot = ASSETS / "homepage-live.png"

    frame_x = T.MARGIN_X
    frame_y = Inches(3.1)
    frame_w = Inches(8.0)
    frame_h = Inches(4.7)

    # Browser chrome shell
    add_round_rect(s, frame_x, frame_y, frame_w, frame_h,
                   fill=T.BG_PANEL_DARK, line=T.DIVIDER_DARK)
    # Top bar
    bar_h = Inches(0.32)
    add_rect(s, frame_x, frame_y, frame_w, bar_h, fill=T.BG_DARK, line=None)
    # Traffic lights
    for i, color in enumerate([T.RED, T.AMBER, T.GREEN]):
        add_oval(s, frame_x + Inches(0.18) + i * Inches(0.22),
                 frame_y + Inches(0.1), Inches(0.13), Inches(0.13),
                 fill=color, line=None)
    # URL pill
    add_round_rect(s, frame_x + Inches(1.3), frame_y + Inches(0.06),
                   Inches(5.2), Inches(0.2), fill=T.BG_PANEL_DARK, line=None)
    add_text(s, frame_x + Inches(1.4), frame_y + Inches(0.04),
             Inches(5.0), Inches(0.24),
             [{"text": "https://product.bookedai.au", "size": T.SZ_TINY,
               "color": T.INK_INV_2, "font": T.FONT_MONO}],
             anchor=MSO_ANCHOR.MIDDLE)

    # Embed screenshot (slightly inset)
    if sshot.exists():
        inset = Inches(0.06)
        img_w = frame_w - inset * 2
        # Compute height assuming 16:9 source
        img_h = Emu(int(int(img_w) * 720 / 1280))
        # Center within frame below the bar
        img_x = frame_x + inset
        img_y = frame_y + bar_h + inset
        # Cap height to frame
        if img_h > (frame_h - bar_h - inset * 2):
            img_h = frame_h - bar_h - inset * 2
        add_image_safe(s, sshot, img_x, img_y, width=img_w)

    # Caption under frame
    add_text(s, frame_x, frame_y + frame_h + Inches(0.05),
             frame_w, Inches(0.3),
             [{"text": C.LIVE_PROOF["screenshot_caption"], "size": T.SZ_TINY,
               "color": T.INK_3, "font": T.FONT_TEXT}], align=PP_ALIGN.CENTER)

    # ----- Right: phone-style Telegram conversation mockup -----
    phone_x = T.MARGIN_X + frame_w + Inches(0.3)
    phone_y = Inches(3.1)
    phone_w = Inches(3.8)
    phone_h = Inches(4.7)
    # Phone shell
    add_round_rect(s, phone_x, phone_y, phone_w, phone_h,
                   fill=T.BG_DARK, line=T.DIVIDER_DARK, line_w=Pt(1.5))
    # Inner screen panel
    inner_inset = Inches(0.12)
    screen_x = phone_x + inner_inset
    screen_y = phone_y + inner_inset
    screen_w = phone_w - inner_inset * 2
    screen_h = phone_h - inner_inset * 2
    add_round_rect(s, screen_x, screen_y, screen_w, screen_h,
                   fill=T.BG_PANEL_DARK, line=None)
    # Top notch / status bar
    bar_h2 = Inches(0.5)
    add_rect(s, screen_x, screen_y, screen_w, bar_h2,
             fill=T.BG_DARK, line=None)
    add_text(s, screen_x + Inches(0.2), screen_y, screen_w - Inches(0.3), bar_h2,
             [{"text": C.LIVE_PROOF["phone_header"], "size": T.SZ_TINY,
               "color": T.ACCENT_SOFT, "bold": True, "font": T.FONT_MONO}],
             anchor=MSO_ANCHOR.MIDDLE)

    # Conversation bubbles
    chat_x = screen_x + Inches(0.18)
    chat_y = screen_y + bar_h2 + Inches(0.18)
    bubble_w_max = screen_w - Inches(0.36)
    cur_y = chat_y
    BUBBLE_COLORS = {
        "customer": (T.BG_PANEL, T.INK),
        "agent": (T.ACCENT, T.INK_INV),
    }
    for role, txt in C.LIVE_PROOF["phone_lines"]:
        bg, fg = BUBBLE_COLORS[role]
        # Estimate bubble dimensions roughly by line count
        n_lines = txt.count("\n") + 1
        # rough width: shorter for long answer-side bubbles, wider for agent
        max_chars = max(len(line) for line in txt.split("\n"))
        approx_w = min(bubble_w_max,
                       Inches(0.18) + Emu(int(min(max_chars, 38) * Inches(0.085) / Inches(1) * 914400)))
        bubble_w = approx_w if approx_w < bubble_w_max else bubble_w_max
        bubble_h = Inches(0.32) + Inches(0.27) * n_lines
        bx = chat_x if role == "customer" else (chat_x + bubble_w_max - bubble_w)
        add_round_rect(s, bx, cur_y, bubble_w, bubble_h,
                       fill=bg, line=None)
        add_text(s, bx + Inches(0.14), cur_y + Inches(0.1),
                 bubble_w - Inches(0.28), bubble_h - Inches(0.1),
                 [{"text": txt, "size": T.SZ_CAPTION, "color": fg,
                   "font": T.FONT_TEXT}])
        cur_y = cur_y + bubble_h + Inches(0.12)

    # ----- Stats strip across bottom (overlaying nothing — left frame) -----
    stat_y = frame_y + frame_h + Inches(0.45)
    add_round_rect(s, frame_x, stat_y, frame_w, Inches(0.55),
                   fill=T.BG_PANEL, line=T.DIVIDER)
    n_stats = len(C.LIVE_PROOF["stats"])
    for i, (val, lbl) in enumerate(C.LIVE_PROOF["stats"]):
        cw = frame_w / n_stats
        x = frame_x + i * cw
        add_text(s, x + Inches(0.2), stat_y, cw - Inches(0.3), Inches(0.55),
                 [
                     {"text": val + "  ", "size": T.SZ_TITLE_S,
                      "color": T.GREEN, "bold": True, "tracking": -10},
                     {"text": lbl, "size": T.SZ_BODY_S,
                      "color": T.INK_2, "font": T.FONT_TEXT},
                 ], anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


# ============================================================
# Slide 05 — Tenant Cases (3-up with hero images)
# ============================================================

def slide_tenant_cases(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Tenant proof cases",
                    "Three live verticals. One reusable template.")

    grid_y = Inches(3.05)
    cell_w = Inches(3.85)
    cell_h = Inches(4.0)
    gap_x = Inches(0.18)
    accents = [T.TENANT_CHESS, T.TENANT_AIMENTOR, T.TENANT_FUTURESWIM]
    for i, case in enumerate(C.TENANT_CASES):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        accent = accents[i]
        add_round_rect(s, x, grid_y, cell_w, cell_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)

        # Hero image / placeholder area
        img_h = Inches(1.7)
        img_inset = Inches(0.12)
        img_x = x + img_inset
        img_y = grid_y + img_inset
        img_w = cell_w - img_inset * 2

        if case["image"] and (ASSETS / case["image"]).exists():
            # Image placeholder background, then embed image fitting width
            add_round_rect(s, img_x, img_y, img_w, img_h,
                           fill=T.BG_LIGHT, line=None)
            # Embed image with width=img_w; height is auto by source ratio
            # source 1200x675 → ratio 16:9 → height = img_w * 9/16
            add_image_safe(s, ASSETS / case["image"], img_x, img_y,
                           width=img_w)
        else:
            # Future Swim — render a custom illustrated panel
            add_round_rect(s, img_x, img_y, img_w, img_h,
                           fill=accent, line=None)
            # 3 horizontal "wave" stripes
            for k in range(3):
                stripe_y = img_y + Inches(0.4) + k * Inches(0.42)
                add_round_rect(s, img_x + Inches(0.3), stripe_y,
                               img_w - Inches(0.6), Inches(0.18),
                               fill=T.INK_INV, line=None)
                # Slightly translucent overlay can't be done easily —
                # leave white stripes against accent for clean swim feel
            add_text(s, img_x, img_y + img_h - Inches(0.6),
                     img_w, Inches(0.5),
                     [{"text": "FUTURE SWIM", "size": T.SZ_BODY,
                       "color": T.INK_INV, "bold": True, "tracking": 220,
                       "font": T.FONT_TEXT}],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Status pill (top-right)
        pill_w = Inches(1.0)
        pill_h = Inches(0.32)
        pill_x = x + cell_w - pill_w - Inches(0.15)
        pill_y = grid_y + Inches(0.18)
        is_live = case["status"] == "LIVE"
        add_round_rect(s, pill_x, pill_y, pill_w, pill_h,
                       fill=T.GREEN if is_live else T.AMBER, line=None)
        add_text(s, pill_x, pill_y, pill_w, pill_h,
                 [{"text": case["status"], "size": T.SZ_TINY,
                   "color": T.INK_INV, "bold": True, "tracking": 200,
                   "font": T.FONT_TEXT}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title + subdomain + body
        body_y = grid_y + img_h + Inches(0.3)
        add_text(s, x + Inches(0.25), body_y, cell_w - Inches(0.4), Inches(0.34),
                 [{"text": case["title"], "size": T.SZ_BODY,
                   "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.25), body_y + Inches(0.36),
                 cell_w - Inches(0.4), Inches(0.3),
                 [{"text": case["subdomain"], "size": T.SZ_CAPTION,
                   "color": accent, "bold": True, "font": T.FONT_MONO}])
        add_text(s, x + Inches(0.25), body_y + Inches(0.7),
                 cell_w - Inches(0.4), Inches(0.55),
                 [{"text": case["tagline"], "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])

        # Loop chip + evidence
        loop_y = body_y + Inches(1.3)
        add_round_rect(s, x + Inches(0.25), loop_y,
                       cell_w - Inches(0.5), Inches(0.4),
                       fill=T.BG_LIGHT, line=T.DIVIDER, line_w=Pt(0.5))
        add_text(s, x + Inches(0.35), loop_y, cell_w - Inches(0.7), Inches(0.4),
                 [{"text": case["loop_short"], "size": T.SZ_TINY,
                   "color": T.INK_2, "font": T.FONT_MONO}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


# ============================================================
# Slide 06 — Market
# ============================================================

def slide_market(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.MARKET["kicker"], C.MARKET["title"])

    # 3 size tiles
    tiles_y = Inches(3.1)
    cell_w = Inches(3.85)
    cell_h = Inches(2.1)
    gap_x = Inches(0.18)
    for i, (label, val, det) in enumerate(C.MARKET["tiles"]):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, tiles_y, cell_w, cell_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.3), tiles_y + Inches(0.18),
                 cell_w - Inches(0.4), Inches(0.3),
                 [{"text": label.upper(), "size": T.SZ_TINY,
                   "color": T.INK_3, "bold": True, "tracking": 220,
                   "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.3), tiles_y + Inches(0.5),
                 cell_w - Inches(0.4), Inches(0.95),
                 [{"text": val, "size": Pt(48), "color": T.ACCENT,
                   "bold": True, "tracking": -30}])
        add_text(s, x + Inches(0.3), tiles_y + Inches(1.4),
                 cell_w - Inches(0.4), cell_h - Inches(1.5),
                 [{"text": det, "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])

    # Build blocks
    bb_y = Inches(5.4)
    add_text(s, T.MARGIN_X, bb_y, Inches(8), Inches(0.3),
             [{"text": "BOTTOM-UP BUILD", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    bullets = []
    for kicker, line in C.MARKET["build_blocks"]:
        bullets.append([
            {"text": kicker + "   ", "size": T.SZ_BODY_S,
             "color": T.ACCENT, "bold": True, "font": T.FONT_MONO},
            {"text": line, "size": T.SZ_BODY_S,
             "color": T.INK_2, "font": T.FONT_TEXT, "space_after": Pt(8)},
        ])
    add_text(s, T.MARGIN_X, bb_y + Inches(0.3),
             Inches(12.13), Inches(1.6), bullets)

    add_slide_number(s, n, total)


# ============================================================
# Slide 07 — Why We Win (quadrant + moats)
# ============================================================

def slide_why_win(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.WHY_WIN["kicker"], C.WHY_WIN["title"])

    # ----- Left: 2D quadrant -----
    plot_x = T.MARGIN_X + Inches(0.6)
    plot_y = Inches(3.1)
    plot_w = Inches(6.6)
    plot_h = Inches(3.9)

    # Plot frame
    add_rect(s, plot_x, plot_y, plot_w, plot_h,
             fill=T.BG_PANEL, line=T.DIVIDER)
    # Mid grid lines
    add_rect(s, plot_x + plot_w / 2, plot_y, Pt(0.75), plot_h,
             fill=T.DIVIDER, line=None)
    add_rect(s, plot_x, plot_y + plot_h / 2, plot_w, Pt(0.75),
             fill=T.DIVIDER, line=None)

    # Quadrant labels in corners
    add_text(s, plot_x + plot_w - Inches(2.2), plot_y + Inches(0.1),
             Inches(2.0), Inches(0.4),
             [{"text": "WHERE WE WIN", "size": T.SZ_TINY,
               "color": T.ACCENT, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}], align=PP_ALIGN.RIGHT)

    # Axis labels
    add_text(s, plot_x, plot_y + plot_h + Inches(0.06),
             plot_w, Inches(0.3),
             [{"text": "VERTICAL DEPTH  →", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 200,
               "font": T.FONT_TEXT}], align=PP_ALIGN.CENTER)
    add_text(s, plot_x - Inches(0.5), plot_y, Inches(0.5), plot_h,
             [{"text": "AI-NATIVE  →", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 200,
               "font": T.FONT_TEXT}], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Plot competitors
    for name, vert, ai in C.WHY_WIN["competitors"]:
        cx = plot_x + Emu(int(vert / 100 * plot_w))
        cy = plot_y + plot_h - Emu(int(ai / 100 * plot_h))
        is_us = name == "BookedAI"
        size = Inches(0.55) if is_us else Inches(0.3)
        color = T.ACCENT if is_us else T.INK_3
        add_oval(s, cx - size / 2, cy - size / 2, size, size,
                 fill=color, line=None)
        # BookedAI ring around it for emphasis
        if is_us:
            ring_size = size + Inches(0.2)
            add_oval(s, cx - ring_size / 2, cy - ring_size / 2,
                     ring_size, ring_size,
                     fill=None, line=T.ACCENT, line_w=Pt(1.5))
        # Label
        lx = cx + Inches(0.32)
        ly = cy - Inches(0.18)
        add_text(s, lx, ly, Inches(2.0), Inches(0.36),
                 [{"text": name, "size": T.SZ_CAPTION,
                   "color": T.ACCENT if is_us else T.INK,
                   "bold": True}])

    # ----- Right: 3 moats -----
    moat_x = plot_x + plot_w + Inches(0.3)
    moat_y = Inches(3.1)
    moat_w = Inches(5.2)
    moat_h = Inches(1.1)
    moat_gap = Inches(0.15)
    add_text(s, moat_x, moat_y - Inches(0.35), moat_w, Inches(0.3),
             [{"text": "THREE MOATS COMPOUND", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    cur_y = moat_y
    for i, (name, body) in enumerate(C.WHY_WIN["moats"]):
        add_round_rect(s, moat_x, cur_y, moat_w, moat_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, moat_x, cur_y, Inches(0.06), moat_h,
                 fill=T.ACCENT, line=None)
        add_text(s, moat_x + Inches(0.25), cur_y + Inches(0.12),
                 moat_w - Inches(0.4), Inches(0.34),
                 [
                     {"text": f"Moat {i+1}  ·  ", "size": T.SZ_TINY,
                      "color": T.ACCENT, "bold": True, "tracking": 200,
                      "font": T.FONT_TEXT},
                     {"text": name, "size": T.SZ_BODY,
                      "color": T.INK, "bold": True},
                 ])
        add_text(s, moat_x + Inches(0.25), cur_y + Inches(0.5),
                 moat_w - Inches(0.4), moat_h - Inches(0.55),
                 [{"text": body, "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])
        cur_y += moat_h + moat_gap

    add_slide_number(s, n, total)


# ============================================================
# Slide 08 — Traction & Roadmap
# ============================================================

def slide_traction(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.TRACTION["kicker"], C.TRACTION["title"])

    # ----- Left: shipped today -----
    col_y = Inches(3.05)
    col_h = Inches(4.0)
    left_w = Inches(5.95)
    right_w = Inches(5.95)
    gap = Inches(0.23)

    add_round_rect(s, T.MARGIN_X, col_y, left_w, col_h,
                   fill=T.BG_PANEL, line=T.DIVIDER)
    add_rect(s, T.MARGIN_X, col_y, left_w, Inches(0.08),
             fill=T.GREEN, line=None)
    add_text(s, T.MARGIN_X + Inches(0.32), col_y + Inches(0.22),
             left_w - Inches(0.5), Inches(0.34),
             [{"text": "SHIPPED  ·  production today", "size": T.SZ_TINY,
               "color": T.GREEN, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    bullets = []
    for line in C.TRACTION["shipped"]:
        bullets.append([
            {"text": "✓  ", "size": T.SZ_BODY_S,
             "color": T.GREEN, "bold": True},
            {"text": line, "size": T.SZ_BODY_S,
             "color": T.INK_2, "font": T.FONT_TEXT, "space_after": Pt(10)},
        ])
    add_text(s, T.MARGIN_X + Inches(0.32), col_y + Inches(0.6),
             left_w - Inches(0.5), col_h - Inches(0.7), bullets)

    # ----- Right: timeline -----
    rx = T.MARGIN_X + left_w + gap
    add_round_rect(s, rx, col_y, right_w, col_h,
                   fill=T.BG_PANEL, line=T.DIVIDER)
    add_rect(s, rx, col_y, right_w, Inches(0.08), fill=T.ACCENT, line=None)
    add_text(s, rx + Inches(0.32), col_y + Inches(0.22),
             right_w - Inches(0.5), Inches(0.34),
             [{"text": "NEXT  ·  Phase 17 → Phase 22+", "size": T.SZ_TINY,
               "color": T.ACCENT, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    # Vertical timeline rail
    rail_x = rx + Inches(0.62)
    rail_top = col_y + Inches(0.85)
    rail_bottom = col_y + col_h - Inches(0.3)
    add_rect(s, rail_x, rail_top, Pt(2), rail_bottom - rail_top,
             fill=T.ACCENT_SOFT, line=None)

    # Phase rows
    n_phases = len(C.TRACTION["phases"])
    rows_y = col_y + Inches(0.85)
    row_h = (col_h - Inches(1.15)) / n_phases
    for i, (phase, dt, descr) in enumerate(C.TRACTION["phases"]):
        y = rows_y + i * row_h
        # Dot
        add_oval(s, rail_x - Inches(0.07), y + Inches(0.1),
                 Inches(0.18), Inches(0.18),
                 fill=T.ACCENT, line=None)
        # Phase + date
        add_text(s, rx + Inches(0.95), y, Inches(2.6), Inches(0.32),
                 [{"text": phase, "size": T.SZ_BODY_S,
                   "color": T.INK, "bold": True}])
        add_text(s, rx + Inches(0.95), y + Inches(0.32), Inches(2.6), Inches(0.28),
                 [{"text": dt, "size": T.SZ_TINY,
                   "color": T.ACCENT, "bold": True, "font": T.FONT_MONO}])
        # Description
        add_text(s, rx + Inches(2.4), y, right_w - Inches(2.5), Inches(0.6),
                 [{"text": descr, "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


# ============================================================
# Slide 09 — Money (business model + unit economics)
# ============================================================

def slide_money(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.MONEY["kicker"], C.MONEY["title"])

    # ----- Top: 3 model components -----
    comp_y = Inches(2.95)
    cell_w = Inches(3.85)
    cell_h = Inches(1.7)
    gap_x = Inches(0.18)
    for i, (name, freq, desc) in enumerate(C.MONEY["components"]):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, comp_y, cell_w, cell_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, x, comp_y, cell_w, Inches(0.07),
                 fill=T.ACCENT, line=None)
        add_text(s, x + Inches(0.28), comp_y + Inches(0.18),
                 cell_w - Inches(0.4), Inches(0.34),
                 [{"text": name, "size": T.SZ_BODY,
                   "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.28), comp_y + Inches(0.52),
                 cell_w - Inches(0.4), Inches(0.3),
                 [{"text": freq.upper(), "size": T.SZ_TINY,
                   "color": T.ACCENT, "bold": True, "tracking": 200,
                   "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.28), comp_y + Inches(0.85),
                 cell_w - Inches(0.4), cell_h - Inches(0.95),
                 [{"text": desc, "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])

    # ----- Bottom: unit economics tiles -----
    ue_y = Inches(4.95)
    add_text(s, T.MARGIN_X, ue_y - Inches(0.3),
             Inches(8), Inches(0.3),
             [{"text": "UNIT ECONOMICS", "size": T.SZ_TINY,
               "color": T.INK_3, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])

    n_ue = len(C.MONEY["unit_econ"])
    ue_w = (Inches(12.13) - gap_x * (n_ue - 1)) / n_ue
    ue_h = Inches(1.85)
    for i, (val, label, sub) in enumerate(C.MONEY["unit_econ"]):
        x = T.MARGIN_X + i * (ue_w + gap_x)
        add_round_rect(s, x, ue_y, ue_w, ue_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.25), ue_y + Inches(0.18),
                 ue_w - Inches(0.4), Inches(0.7),
                 [{"text": val, "size": T.SZ_METRIC_S,
                   "color": T.ACCENT, "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.25), ue_y + Inches(0.85),
                 ue_w - Inches(0.4), Inches(0.32),
                 [{"text": label, "size": T.SZ_BODY_S,
                   "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.25), ue_y + Inches(1.2),
                 ue_w - Inches(0.4), ue_h - Inches(1.3),
                 [{"text": sub, "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


# ============================================================
# Slide 10 — Team (4 portraits from pitch.bookedai.au)
# ============================================================

def slide_team(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.TEAM_KICKER, C.TEAM_TITLE)

    # Sub-body under header
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(12.13), Inches(0.55),
             [{"text": C.TEAM_BODY, "size": T.SZ_BODY_S,
               "color": T.INK_2, "font": T.FONT_TEXT}])

    # Four-up grid
    grid_y = Inches(3.65)
    cell_w = Inches(2.95)
    cell_h = Inches(3.4)
    gap_x = Inches(0.18)
    for i, m in enumerate(C.TEAM_MEMBERS):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, grid_y, cell_w, cell_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)

        # Portrait area — 16:9 banner crop
        img_inset = Inches(0.12)
        img_x = x + img_inset
        img_y = grid_y + img_inset
        img_w = cell_w - img_inset * 2
        # Source ratio is 900x502 (~16:9) → height auto from width
        img_path = ASSETS / m["image"]
        if img_path.exists():
            # Background panel in case the image has alpha gaps
            add_round_rect(s, img_x, img_y, img_w, Inches(1.45),
                           fill=T.BG_DARK_2, line=None)
            add_image_safe(s, img_path, img_x, img_y, width=img_w)

        # Name + role
        body_y = grid_y + Inches(1.7)
        add_text(s, x + Inches(0.22), body_y, cell_w - Inches(0.4), Inches(0.36),
                 [{"text": m["name"], "size": T.SZ_BODY,
                   "color": T.INK, "bold": True}])

        # Role pill
        pill_w = Inches(0.85)
        pill_h = Inches(0.28)
        pill_y = body_y + Inches(0.4)
        add_round_rect(s, x + Inches(0.22), pill_y, pill_w, pill_h,
                       fill=T.ACCENT, line=None)
        add_text(s, x + Inches(0.22), pill_y, pill_w, pill_h,
                 [{"text": m["role"], "size": T.SZ_TINY,
                   "color": T.INK_INV, "bold": True, "tracking": 200,
                   "font": T.FONT_TEXT}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Badges (small text, comma separated)
        badges_y = pill_y + Inches(0.36)
        badges_text = "  ·  ".join(m["badges"])
        add_text(s, x + Inches(0.22), badges_y, cell_w - Inches(0.4), Inches(0.3),
                 [{"text": badges_text, "size": T.SZ_TINY,
                   "color": T.INK_3, "bold": True, "tracking": 100,
                   "font": T.FONT_MONO}])

        # Bio
        bio_y = badges_y + Inches(0.34)
        add_text(s, x + Inches(0.22), bio_y, cell_w - Inches(0.4),
                 cell_h - Inches(2.5),
                 [{"text": m["bio"], "size": T.SZ_CAPTION,
                   "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


# ============================================================
# Slide 11 — The Ask (use of funds + milestones)
# ============================================================

def slide_ask(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, T.BG_DARK)
    add_brand_mark(s, dark=True)

    add_kicker(s, T.MARGIN_X, Inches(1.35), C.ASK["kicker"], color=T.ACCENT_SOFT)
    add_text(s, T.MARGIN_X, Inches(1.7), Inches(11), Inches(1.1),
             [{"text": C.ASK["headline"], "size": Pt(58),
               "color": T.INK_INV, "bold": True, "tracking": -30}])
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(11), Inches(0.5),
             [{"text": C.ASK["structure"], "size": T.SZ_SUBTITLE,
               "color": T.ACCENT_SOFT, "font": T.FONT_TEXT}])

    # Use of funds — 4 columns (taller now since we have more vertical room)
    uof_y = Inches(3.7)
    cell_w = Inches(2.9)
    cell_h = Inches(1.85)
    gap_x = Inches(0.16)
    for i, (pct, label, descr) in enumerate(C.ASK["use_of_funds"]):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, uof_y, cell_w, cell_h,
                       fill=T.BG_PANEL_DARK, line=T.DIVIDER_DARK)
        add_text(s, x + Inches(0.22), uof_y + Inches(0.16),
                 cell_w - Inches(0.4), Inches(0.7),
                 [{"text": pct, "size": Pt(40),
                   "color": T.ACCENT_SOFT, "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.22), uof_y + Inches(0.85),
                 cell_w - Inches(0.4), Inches(0.34),
                 [{"text": label, "size": T.SZ_BODY,
                   "color": T.INK_INV, "bold": True}])
        add_text(s, x + Inches(0.22), uof_y + Inches(1.2),
                 cell_w - Inches(0.4), cell_h - Inches(1.3),
                 [{"text": descr, "size": T.SZ_CAPTION,
                   "color": T.INK_INV_2, "font": T.FONT_TEXT}])

    # Milestones strip
    ms_y = Inches(5.85)
    add_text(s, T.MARGIN_X, ms_y, Inches(8), Inches(0.3),
             [{"text": "MILESTONES THIS ROUND FUNDS",
               "size": T.SZ_TINY, "color": T.ACCENT_SOFT,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for line in C.ASK["milestones"]:
        bullets.append([
            {"text": "→  ", "size": T.SZ_BODY_S,
             "color": T.ACCENT_SOFT, "bold": True},
            {"text": line, "size": T.SZ_BODY_S,
             "color": T.INK_INV, "font": T.FONT_TEXT, "space_after": Pt(4)},
        ])
    add_text(s, T.MARGIN_X, ms_y + Inches(0.3),
             Inches(12.13), Inches(1.4), bullets)

    add_slide_number(s, n, total, dark=True)


# ============================================================
# Slide 12 — Thank You / CTA (final, with hero image)
# ============================================================

def slide_thank_you(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, T.BG_DARK)

    # ----- Right: hero image -----
    img_path = ASSETS / C.THANK_YOU["image"]
    img_w = Inches(6.0)
    img_x = T.SLIDE_WIDTH - T.MARGIN_X - img_w
    img_y = Inches(0.7)
    img_h = Inches(6.1)

    if img_path.exists():
        # Backing panel for clean edges
        add_round_rect(s, img_x, img_y, img_w, img_h,
                       fill=T.BG_PANEL_DARK, line=T.DIVIDER_DARK, line_w=Pt(1))
        # Insert image; rely on natural aspect (1400x933 → 3:2)
        # Width-fit; height auto. Source aspect → height ≈ img_w * 933/1400 ≈ 4.0"
        # We'll center it vertically within the panel.
        natural_h = Emu(int(int(img_w) * 933 / 1400))
        if natural_h > img_h:
            natural_h = img_h
        image_top = img_y + (img_h - natural_h) / 2
        add_image_safe(s, img_path, img_x, image_top, width=img_w)
        # Subtle overlay band at bottom for caption
        cap_h = Inches(0.4)
        add_rect(s, img_x, img_y + img_h - cap_h, img_w, cap_h,
                 fill=T.BG_DARK, line=None)
        add_text(s, img_x + Inches(0.25), img_y + img_h - cap_h,
                 img_w - Inches(0.4), cap_h,
                 [{"text": "Live customer-care continuity proof on bookedai.au",
                   "size": T.SZ_TINY, "color": T.INK_INV_2,
                   "font": T.FONT_TEXT}], anchor=MSO_ANCHOR.MIDDLE)

    # ----- Left: text + CTAs -----
    text_x = T.MARGIN_X
    text_w = img_x - T.MARGIN_X - Inches(0.4)

    # Brand mark top-left
    add_brand_mark(s, dark=True, t=Inches(0.55), h=Inches(0.5))

    # Hero text
    add_text(s, text_x, Inches(1.7), text_w, Inches(1.4),
             [{"text": C.THANK_YOU["headline"], "size": Pt(80),
               "color": T.INK_INV, "bold": True, "tracking": -50}])
    add_text(s, text_x, Inches(3.0), text_w, Inches(0.7),
             [{"text": C.THANK_YOU["subhead"], "size": Pt(24),
               "color": T.ACCENT_SOFT, "bold": True, "tracking": -10,
               "font": T.FONT_DISPLAY}])

    # Hairline accent
    add_rect(s, text_x, Inches(3.85), Inches(1.5), Pt(2),
             fill=T.ACCENT, line=None)

    # CTA stack — 4 chips vertical on the left
    cta_y0 = Inches(4.1)
    chip_h = Inches(0.55)
    chip_gap = Inches(0.12)
    for i, (lbl, val) in enumerate(C.THANK_YOU["ctas"]):
        cy = cta_y0 + i * (chip_h + chip_gap)
        add_round_rect(s, text_x, cy, text_w, chip_h,
                       fill=T.BG_PANEL_DARK, line=T.DIVIDER_DARK)
        # Accent dot
        add_oval(s, text_x + Inches(0.18), cy + Inches(0.18),
                 Inches(0.18), Inches(0.18),
                 fill=T.ACCENT, line=None)
        # Label + value
        add_text(s, text_x + Inches(0.5), cy, text_w - Inches(0.6), chip_h,
                 [
                     {"text": lbl + "  ", "size": T.SZ_TINY,
                      "color": T.INK_INV_3, "bold": True, "tracking": 200,
                      "font": T.FONT_TEXT},
                     {"text": val, "size": T.SZ_BODY_S,
                      "color": T.INK_INV, "bold": True, "font": T.FONT_MONO},
                 ], anchor=MSO_ANCHOR.MIDDLE)

    # Closer line at bottom
    add_text(s, T.MARGIN_X, T.SLIDE_HEIGHT - Inches(0.55),
             T.SLIDE_WIDTH - T.MARGIN_X * 2, Inches(0.32),
             [{"text": C.THANK_YOU["closer"], "size": T.SZ_CAPTION,
               "color": T.INK_INV_3, "font": T.FONT_TEXT}],
             align=PP_ALIGN.CENTER)

    add_slide_number(s, n, total, dark=True)


# ============================================================
# Build orchestration
# ============================================================

SLIDES = [
    ("01_cover",         slide_cover),
    ("02_problem",       slide_problem),
    ("03_solution",      slide_solution),
    ("04_live_proof",    slide_live_proof),
    ("05_tenant_cases",  slide_tenant_cases),
    ("06_market",        slide_market),
    ("07_why_win",       slide_why_win),
    ("08_traction",      slide_traction),
    ("09_money",         slide_money),
    ("10_team",          slide_team),
    ("11_ask",           slide_ask),
    ("12_thank_you",     slide_thank_you),
]


def build():
    prs = Presentation()
    prs.slide_width = T.SLIDE_WIDTH
    prs.slide_height = T.SLIDE_HEIGHT

    total = len(SLIDES)
    for i, (name, fn) in enumerate(SLIDES, start=1):
        fn(prs, i, total)

    OUTPUT.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT / f"BookedAI_5min_Pitch_{C.DECK_DATE.isoformat()}.pptx"
    prs.save(str(out_path))
    size_kb = out_path.stat().st_size / 1024
    print(f"✓ Built {out_path.name}  ·  {total} slides  ·  {size_kb:,.0f} KB")
    print(f"  → {out_path}")
    return out_path


if __name__ == "__main__":
    build()
