#!/usr/bin/env python3
"""
BookedAI investor pitch deck builder.

Produces output/BookedAI_Investor_Pitch_<YYYY-MM-DD>.pptx
in 16:9, Apple-inspired design system locked in DESIGN.md.

Run:
    cd pitch-deck
    ../.venv-pdf/bin/python build_pitch_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

import content as C
import theme as T


HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUTPUT = HERE / "output"


# ---------- Low-level helpers ----------

def _set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_no_fill(shape):
    shape.fill.background()


def _set_line(shape, color: RGBColor | None, width=Pt(0.75)):
    line = shape.line
    if color is None:
        line.fill.background()
        return
    line.color.rgb = color
    line.width = width


def add_rect(slide, left, top, width, height, *, fill=None, line=None, line_w=Pt(0.75), shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.shadow.inherit = False
    if fill is None:
        _set_no_fill(s)
    else:
        _set_fill(s, fill)
    _set_line(s, line, line_w)
    return s


def add_round_rect(slide, left, top, width, height, *, fill=None, line=None, line_w=Pt(0.75)):
    return add_rect(slide, left, top, width, height,
                    fill=fill, line=line, line_w=line_w,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def add_text(slide, left, top, width, height, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None,
             margin=Inches(0.0)):
    """
    runs: list of dicts: {text, size, color, bold?, font?, tracking?}
          or list of lists for multiple paragraphs.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        _set_fill(box, fill)
    if line is not None:
        _set_line(box, line)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    tf.vertical_anchor = anchor

    paragraphs = runs if (runs and isinstance(runs[0], list)) else [runs]
    for p_index, paragraph in enumerate(paragraphs):
        if p_index == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        for r_index, r in enumerate(paragraph):
            if r_index == 0 and p_index == 0:
                run = para.runs[0] if para.runs else para.add_run()
            else:
                run = para.add_run()
            run.text = r.get("text", "")
            T.apply_text_style(
                run,
                size=r.get("size", T.SZ_BODY),
                color=r.get("color", T.INK),
                bold=r.get("bold", False),
                font=r.get("font", T.FONT_DISPLAY),
                tracking=r.get("tracking"),
            )
        if "space_after" in (paragraph[-1] if paragraph else {}):
            para.space_after = paragraph[-1]["space_after"]
    return box


def add_kicker(slide, left, top, text, color=T.ACCENT, *, width=Inches(6)):
    return add_text(slide, left, top, width, Inches(0.32),
                    [{"text": text.upper(), "size": T.SZ_KICKER, "color": color,
                      "bold": True, "font": T.FONT_TEXT, "tracking": 240}])


def add_title(slide, left, top, text, *, width=Inches(11.5), color=T.INK, size=T.SZ_TITLE, bold=True):
    return add_text(slide, left, top, width, Inches(1.4),
                    [{"text": text, "size": size, "color": color, "bold": bold,
                      "font": T.FONT_DISPLAY, "tracking": -20}])


def add_footer(slide, *, dark=False):
    color = T.INK_INV_3 if dark else T.INK_3
    add_text(slide, T.MARGIN_X, T.SLIDE_HEIGHT - Inches(0.32),
             Inches(8), Inches(0.22),
             [{"text": f"BookedAI · {C.URLS['homepage']} · Confidential", "size": T.SZ_TINY,
               "color": color, "font": T.FONT_TEXT}])
    pn_text = "" if dark else ""
    # Slide number drawn separately by builder loop


def add_slide_number(slide, n, total, *, dark=False):
    color = T.INK_INV_3 if dark else T.INK_3
    add_text(slide, T.SLIDE_WIDTH - Inches(1.2), T.SLIDE_HEIGHT - Inches(0.32),
             Inches(0.9), Inches(0.22),
             [{"text": f"{n:02d} / {total:02d}", "size": T.SZ_TINY, "color": color,
               "font": T.FONT_TEXT}], align=PP_ALIGN.RIGHT)


def add_image_safe(slide, path: Path, left, top, *, width=None, height=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def fill_slide_bg(slide, color: RGBColor):
    """Fill the entire slide background (using a backing rectangle)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, T.SLIDE_WIDTH, T.SLIDE_HEIGHT)
    _set_fill(s, color)
    _set_line(s, None)
    s.shadow.inherit = False
    # Send to back
    spTree = s._element.getparent()
    spTree.remove(s._element)
    spTree.insert(2, s._element)
    return s


def add_brand_mark(slide, *, left=None, top=Inches(0.45), height=Inches(0.42), dark=False):
    """Tiny BookedAI brand mark in the corner."""
    if left is None:
        left = T.MARGIN_X
    logo = ASSETS / ("logo-light.png" if dark else "logo-dark.png")
    if logo.exists():
        add_image_safe(slide, logo, left, top, height=height)
    else:
        color = T.INK_INV if dark else T.INK
        add_text(slide, left, top, Inches(2), height,
                 [{"text": "BookedAI", "size": T.SZ_SUBTITLE, "color": color, "bold": True}])


# ---------- Slide builders ----------

def slide_cover(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_slide_bg(s, T.BG_DARK)

    # Decorative gradient mark right side
    mark = ASSETS / "mark-gradient.png"
    if mark.exists():
        add_image_safe(s, mark, T.SLIDE_WIDTH - Inches(4.8), Inches(1.0),
                       height=Inches(5.5))

    # Big logo top-left
    add_brand_mark(s, top=Inches(0.6), height=Inches(0.55), dark=True)

    # Hero
    add_text(s, T.MARGIN_X, Inches(2.4), Inches(9), Inches(1.0),
             [{"text": C.COMPANY, "size": T.SZ_HERO, "color": T.INK_INV,
               "bold": True, "tracking": -40}])
    add_text(s, T.MARGIN_X, Inches(3.6), Inches(10), Inches(0.7),
             [{"text": C.TAGLINE, "size": Pt(34), "color": T.ACCENT_SOFT,
               "bold": True, "tracking": -20}])
    add_text(s, T.MARGIN_X, Inches(4.5), Inches(10), Inches(0.5),
             [{"text": C.SUBTAGLINE, "size": T.SZ_SUBTITLE, "color": T.INK_INV_2,
               "font": T.FONT_TEXT}])

    # Hairline
    add_rect(s, T.MARGIN_X, Inches(5.6), Inches(1.5), Pt(2), fill=T.ACCENT, line=None)

    # Footer block
    add_text(s, T.MARGIN_X, Inches(5.85), Inches(7), Inches(0.34),
             [{"text": C.AUDIENCE, "size": T.SZ_BODY_S, "color": T.INK_INV_2,
               "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X, Inches(6.2), Inches(7), Inches(0.34),
             [{"text": C.PRESENTER, "size": T.SZ_BODY_S, "color": T.INK_INV_2,
               "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X, Inches(6.55), Inches(7), Inches(0.34),
             [{"text": C.DECK_DATE.isoformat(), "size": T.SZ_CAPTION, "color": T.INK_INV_3,
               "font": T.FONT_TEXT}])
    add_slide_number(s, n, total, dark=True)


def _content_header(s, kicker, title, *, dark=False):
    fill_slide_bg(s, T.BG_DARK if dark else T.BG_LIGHT)
    add_brand_mark(s, dark=dark)
    add_kicker(s, T.MARGIN_X, Inches(1.35), kicker)
    add_title(s, T.MARGIN_X, Inches(1.7), title,
              color=T.INK_INV if dark else T.INK,
              size=T.SZ_TITLE)


def slide_vision(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.VISION["kicker"], C.VISION["title"])

    # Body
    add_text(s, T.MARGIN_X, Inches(3.4), Inches(11.5), Inches(1.4),
             [{"text": C.VISION["body"], "size": T.SZ_SUBTITLE, "color": T.INK_2,
               "font": T.FONT_TEXT}])

    # Three pillars
    pillar_y = Inches(5.2)
    pillar_w = Inches(3.85)
    pillar_h = Inches(1.6)
    gap = Inches(0.18)
    for i, (k, v) in enumerate(C.VISION["pillars"]):
        x = T.MARGIN_X + i * (pillar_w + gap)
        add_round_rect(s, x, pillar_y, pillar_w, pillar_h,
                       fill=T.BG_PANEL, line=T.DIVIDER, line_w=Pt(0.75))
        add_text(s, x + Inches(0.32), pillar_y + Inches(0.22), pillar_w - Inches(0.4), Inches(0.5),
                 [{"text": k, "size": T.SZ_TITLE_S, "color": T.ACCENT, "bold": True}])
        add_text(s, x + Inches(0.32), pillar_y + Inches(0.85), pillar_w - Inches(0.4), Inches(0.7),
                 [{"text": v, "size": T.SZ_BODY, "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.PROBLEM["kicker"], C.PROBLEM["title"])

    # Stats row
    stats_y = Inches(3.3)
    stats_w = Inches(3.85)
    stats_h = Inches(1.5)
    gap = Inches(0.18)
    for i, (val, label) in enumerate(C.PROBLEM["stats"]):
        x = T.MARGIN_X + i * (stats_w + gap)
        add_round_rect(s, x, stats_y, stats_w, stats_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.32), stats_y + Inches(0.22), stats_w - Inches(0.4), Inches(0.7),
                 [{"text": val, "size": T.SZ_METRIC_S, "color": T.RED, "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.32), stats_y + Inches(0.85), stats_w - Inches(0.4), Inches(0.6),
                 [{"text": label, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Narrative bullets
    narr_y = Inches(5.1)
    runs = []
    for line in C.PROBLEM["narrative"]:
        runs.append([
            {"text": "→ ", "size": T.SZ_BODY, "color": T.ACCENT, "bold": True},
            {"text": line, "size": T.SZ_BODY, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(8)},
        ])
    add_text(s, T.MARGIN_X, narr_y, Inches(11.8), Inches(2), runs)

    add_slide_number(s, n, total)


def slide_why_now(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Why Now", "The market window is open — for the first time, the economics work.")

    # Stats tiles
    stats_y = Inches(3.2)
    stats_w = Inches(3.85)
    stats_h = Inches(1.7)
    gap = Inches(0.18)
    for i, (val, label, sub) in enumerate(C.WHY_NOW_STATS):
        x = T.MARGIN_X + i * (stats_w + gap)
        add_round_rect(s, x, stats_y, stats_w, stats_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.32), stats_y + Inches(0.20), stats_w - Inches(0.4), Inches(0.7),
                 [{"text": val, "size": T.SZ_METRIC_S, "color": T.ACCENT, "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.32), stats_y + Inches(0.78), stats_w - Inches(0.4), Inches(0.32),
                 [{"text": label, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.32), stats_y + Inches(1.10), stats_w - Inches(0.4), Inches(0.55),
                 [{"text": sub, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Signals
    sig_y = Inches(5.15)
    sig_w = Inches(3.85)
    sig_h = Inches(1.65)
    for i, (title, body) in enumerate(C.WHY_NOW_SIGNALS):
        x = T.MARGIN_X + i * (sig_w + gap)
        add_text(s, x, sig_y, sig_w, Inches(0.34),
                 [{"text": title, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x, sig_y + Inches(0.34), sig_w, sig_h - Inches(0.3),
                 [{"text": body, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_solution(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.SOLUTION["kicker"], C.SOLUTION["title"])

    # Loop diagram
    steps = C.SOLUTION["loop_steps"]
    loop_y = Inches(3.45)
    pill_h = Inches(0.5)
    total_w = Inches(12.13)
    n_steps = len(steps)
    n_arrows = n_steps - 1
    arrow_w = Inches(0.25)
    pill_w_emu = (total_w - n_arrows * arrow_w) / n_steps
    cur = T.MARGIN_X
    for i, step in enumerate(steps):
        add_round_rect(s, cur, loop_y, pill_w_emu, pill_h,
                       fill=T.ACCENT if i in (3, 5) else T.BG_PANEL,
                       line=T.ACCENT if i in (3, 5) else T.DIVIDER)
        add_text(s, cur, loop_y, pill_w_emu, pill_h,
                 [{"text": step, "size": T.SZ_BODY, "bold": True,
                   "color": T.INK_INV if i in (3, 5) else T.INK}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur = cur + pill_w_emu
        if i < n_arrows:
            add_text(s, cur, loop_y, arrow_w, pill_h,
                     [{"text": "›", "size": T.SZ_TITLE_S, "color": T.INK_3, "bold": True}],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cur = cur + arrow_w

    add_text(s, T.MARGIN_X, loop_y + Inches(0.65), Inches(12.13), Inches(0.5),
             [{"text": C.SOLUTION["loop_caption"], "size": T.SZ_BODY_S,
               "color": T.INK_3, "font": T.FONT_TEXT}], align=PP_ALIGN.CENTER)

    # Three agent classes
    agent_y = Inches(5.0)
    agent_w = Inches(3.85)
    agent_h = Inches(1.85)
    gap = Inches(0.18)
    for i, (title, body) in enumerate(C.SOLUTION["agent_classes"]):
        x = T.MARGIN_X + i * (agent_w + gap)
        add_round_rect(s, x, agent_y, agent_w, agent_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, x, agent_y, Inches(0.06), agent_h, fill=T.ACCENT, line=None)
        add_text(s, x + Inches(0.28), agent_y + Inches(0.22), agent_w - Inches(0.4), Inches(0.6),
                 [{"text": title, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.28), agent_y + Inches(0.7), agent_w - Inches(0.4), agent_h - Inches(0.7),
                 [{"text": body, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_surfaces(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Product surfaces", "Twelve subdomains. One operating system.")

    # 4-column grid of surface cards
    grid_y = Inches(3.1)
    cols = 4
    rows = 3
    cell_w = Inches(3.0)
    cell_h = Inches(1.2)
    gap_x = Inches(0.13)
    gap_y = Inches(0.15)
    for idx, (sub, descr) in enumerate(C.SURFACES[: cols * rows]):
        r, c = divmod(idx, cols)
        x = T.MARGIN_X + c * (cell_w + gap_x)
        y = grid_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.18), y + Inches(0.14), cell_w - Inches(0.3), Inches(0.34),
                 [{"text": sub, "size": T.SZ_BODY_S, "color": T.ACCENT, "bold": True,
                   "font": T.FONT_MONO}])
        add_text(s, x + Inches(0.18), y + Inches(0.5), cell_w - Inches(0.3), cell_h - Inches(0.55),
                 [{"text": descr, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_journey(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.JOURNEY["kicker"], C.JOURNEY["title"])

    # 7 columns one row
    grid_y = Inches(3.4)
    cell_w = Inches(1.72)
    cell_h = Inches(2.6)
    gap_x = Inches(0.07)
    cur_x = T.MARGIN_X
    for i, (step, desc) in enumerate(C.JOURNEY["steps"]):
        # Numbered circle
        add_round_rect(s, cur_x + Inches(0.6), grid_y, Inches(0.5), Inches(0.5),
                       fill=T.ACCENT, line=None)
        add_text(s, cur_x + Inches(0.6), grid_y, Inches(0.5), Inches(0.5),
                 [{"text": str(i + 1), "size": T.SZ_BODY, "color": T.INK_INV, "bold": True}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Step label
        add_text(s, cur_x, grid_y + Inches(0.65), cell_w, Inches(0.4),
                 [{"text": step, "size": T.SZ_BODY, "color": T.INK, "bold": True}],
                 align=PP_ALIGN.CENTER)
        # Description
        add_text(s, cur_x + Inches(0.05), grid_y + Inches(1.05), cell_w - Inches(0.1), cell_h - Inches(1.1),
                 [{"text": desc, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}],
                 align=PP_ALIGN.CENTER)
        cur_x = cur_x + cell_w + gap_x

    # Footer
    add_text(s, T.MARGIN_X, Inches(6.4), Inches(12.13), Inches(0.32),
             [{"text": C.JOURNEY["footer"], "size": T.SZ_CAPTION, "color": T.INK_3,
               "font": T.FONT_TEXT}], align=PP_ALIGN.CENTER)

    add_slide_number(s, n, total)


def slide_agents(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Three Agent Classes", "Search, Revenue Ops, Customer Care — all production-grade.")

    # 3 column cards
    grid_y = Inches(3.2)
    cell_w = Inches(3.85)
    cell_h = Inches(3.6)
    gap_x = Inches(0.18)
    for i, agent in enumerate(C.AGENTS):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, grid_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        # Header strip
        add_round_rect(s, x, grid_y, cell_w, Inches(0.85), fill=T.ACCENT_DEEP, line=None)
        add_text(s, x + Inches(0.32), grid_y + Inches(0.14), cell_w - Inches(0.5), Inches(0.36),
                 [{"text": agent["title"], "size": T.SZ_BODY, "color": T.INK_INV, "bold": True}])
        add_text(s, x + Inches(0.32), grid_y + Inches(0.5), cell_w - Inches(0.5), Inches(0.32),
                 [{"text": agent["subtitle"], "size": T.SZ_CAPTION, "color": T.INK_INV_2,
                   "font": T.FONT_TEXT}])
        # Channels
        add_text(s, x + Inches(0.32), grid_y + Inches(1.05), cell_w - Inches(0.5), Inches(0.4),
                 [{"text": "CHANNELS", "size": T.SZ_TINY, "color": T.INK_3,
                   "bold": True, "tracking": 200, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.32), grid_y + Inches(1.3), cell_w - Inches(0.5), Inches(0.6),
                 [{"text": agent["channels"], "size": T.SZ_BODY_S, "color": T.INK_2,
                   "font": T.FONT_TEXT}])
        # Responsibilities
        add_text(s, x + Inches(0.32), grid_y + Inches(1.85), cell_w - Inches(0.5), Inches(0.4),
                 [{"text": "RESPONSIBILITIES", "size": T.SZ_TINY, "color": T.INK_3,
                   "bold": True, "tracking": 200, "font": T.FONT_TEXT}])
        bullets = []
        for line in agent["responsibilities"]:
            bullets.append([
                {"text": "•  ", "size": T.SZ_CAPTION, "color": T.ACCENT, "bold": True},
                {"text": line, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT,
                 "space_after": Pt(6)},
            ])
        add_text(s, x + Inches(0.32), grid_y + Inches(2.15), cell_w - Inches(0.5), Inches(1.4), bullets)

    add_slide_number(s, n, total)


def slide_live_evidence(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Live evidence", "Telegram → Admin → Workspace, in under 30 seconds.")

    grid_y = Inches(3.1)
    cell_w = Inches(3.85)
    cell_h = Inches(3.7)
    gap_x = Inches(0.18)
    SPEAKER_COLORS = {
        "customer": T.ACCENT,
        "agent": T.GREEN,
        "system": T.INK_3,
    }
    for i, frame in enumerate(C.LIVE_EVIDENCE):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, grid_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        # Step + surface header
        add_text(s, x + Inches(0.28), grid_y + Inches(0.18), cell_w - Inches(0.4), Inches(0.32),
                 [{"text": f"STEP {frame['step']}  ·  {frame['surface']}",
                   "size": T.SZ_TINY, "color": T.ACCENT, "bold": True, "tracking": 180,
                   "font": T.FONT_MONO}])
        # Title
        add_text(s, x + Inches(0.28), grid_y + Inches(0.48), cell_w - Inches(0.4), Inches(0.7),
                 [{"text": frame["title"], "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        # Caption
        add_text(s, x + Inches(0.28), grid_y + Inches(1.18), cell_w - Inches(0.4), Inches(0.6),
                 [{"text": frame["caption"], "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])
        # Conversation
        cur_y = grid_y + Inches(1.85)
        for spk, txt in frame["lines"]:
            color = SPEAKER_COLORS.get(spk, T.INK_3)
            label = spk.upper()
            add_text(s, x + Inches(0.28), cur_y, cell_w - Inches(0.4), Inches(0.55),
                     [
                         {"text": f"{label}  ", "size": T.SZ_TINY, "color": color, "bold": True,
                          "tracking": 160, "font": T.FONT_MONO},
                         {"text": txt, "size": T.SZ_CAPTION, "color": T.INK, "font": T.FONT_TEXT},
                     ])
            cur_y = cur_y + Inches(0.42)

    add_slide_number(s, n, total)


def slide_tenant_case(prs, n, total, case, accent):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, f"Tenant case · {case['name']}", case["title"])

    # Subdomain tag
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(7), Inches(0.3),
             [{"text": case["subdomain"], "size": T.SZ_BODY_S, "color": accent,
               "bold": True, "font": T.FONT_MONO}])

    # Tagline
    add_text(s, T.MARGIN_X, Inches(3.3), Inches(11.5), Inches(0.55),
             [{"text": case["tagline"], "size": T.SZ_SUBTITLE, "color": T.INK_2,
               "font": T.FONT_TEXT}])

    # Loop box
    loop_y = Inches(4.05)
    add_round_rect(s, T.MARGIN_X, loop_y, Inches(12.13), Inches(0.95),
                   fill=T.BG_PANEL, line=T.DIVIDER)
    add_text(s, T.MARGIN_X + Inches(0.3), loop_y + Inches(0.12), Inches(2.5), Inches(0.3),
             [{"text": "LOOP", "size": T.SZ_TINY, "color": accent, "bold": True,
               "tracking": 220, "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X + Inches(0.3), loop_y + Inches(0.4), Inches(11.6), Inches(0.5),
             [{"text": case["loop"], "size": T.SZ_BODY_S, "color": T.INK,
               "font": T.FONT_TEXT}])

    # Evidence bullets
    ev_y = Inches(5.2)
    add_text(s, T.MARGIN_X, ev_y, Inches(8), Inches(0.3),
             [{"text": "EVIDENCE", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for line in case["evidence"]:
        bullets.append([
            {"text": "•  ", "size": T.SZ_BODY, "color": accent, "bold": True},
            {"text": line, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(8)},
        ])
    add_text(s, T.MARGIN_X, ev_y + Inches(0.32), Inches(8), Inches(2), bullets)

    # Status badge
    badge_w = Inches(4.0)
    badge_h = Inches(2.0)
    badge_x = T.SLIDE_WIDTH - T.MARGIN_X - badge_w
    badge_y = ev_y
    add_round_rect(s, badge_x, badge_y, badge_w, badge_h, fill=accent, line=None)
    add_text(s, badge_x + Inches(0.3), badge_y + Inches(0.32), badge_w - Inches(0.5), Inches(0.4),
             [{"text": "STATUS", "size": T.SZ_TINY, "color": T.INK_INV_2, "bold": True,
               "tracking": 220, "font": T.FONT_TEXT}])
    add_text(s, badge_x + Inches(0.3), badge_y + Inches(0.7), badge_w - Inches(0.5), badge_h - Inches(0.8),
             [{"text": case["status"], "size": T.SZ_TITLE_S, "color": T.INK_INV,
               "bold": True, "tracking": -10}])

    add_slide_number(s, n, total)


def slide_architecture(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Architecture", "Six layers · production on a single VPS · self-hosted by design.")

    rows_y = Inches(3.2)
    row_h = Inches(0.62)
    row_gap = Inches(0.08)
    for i, (layer, path, desc) in enumerate(C.ARCHITECTURE_LAYERS):
        y = rows_y + i * (row_h + row_gap)
        add_round_rect(s, T.MARGIN_X, y, Inches(12.13), row_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, T.MARGIN_X, y, Inches(0.08), row_h, fill=T.ACCENT, line=None)
        add_text(s, T.MARGIN_X + Inches(0.3), y, Inches(2.6), row_h,
                 [{"text": layer, "size": T.SZ_BODY, "color": T.INK, "bold": True}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(2.9), y, Inches(2.5), row_h,
                 [{"text": path, "size": T.SZ_BODY_S, "color": T.ACCENT,
                   "bold": True, "font": T.FONT_MONO}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(5.4), y, Inches(6.6), row_h,
                 [{"text": desc, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


def slide_tech_stack(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Tech Stack", "Boring, proven, fast — purposefully unsexy.")

    layout = [
        ("Frontend", C.TECH_STACK["frontend"]),
        ("Backend", C.TECH_STACK["backend"]),
        ("Data", C.TECH_STACK["data"]),
        ("Automation", C.TECH_STACK["automation"]),
        ("Infrastructure", C.TECH_STACK["infra"]),
        ("AI / ML", C.TECH_STACK["ai"]),
    ]
    grid_y = Inches(3.15)
    cell_w = Inches(4.0)
    cell_h = Inches(1.8)
    gap_x = Inches(0.1)
    gap_y = Inches(0.1)
    for idx, (head, items) in enumerate(layout):
        r, c = divmod(idx, 3)
        x = T.MARGIN_X + c * (cell_w + gap_x)
        y = grid_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.28), y + Inches(0.18), cell_w - Inches(0.4), Inches(0.34),
                 [{"text": head, "size": T.SZ_BODY, "color": T.ACCENT, "bold": True}])
        bullets = []
        for it in items:
            bullets.append([
                {"text": "›  ", "size": T.SZ_CAPTION, "color": T.INK_3, "bold": True,
                 "font": T.FONT_MONO},
                {"text": it, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT,
                 "space_after": Pt(3)},
            ])
        add_text(s, x + Inches(0.28), y + Inches(0.52), cell_w - Inches(0.4), cell_h - Inches(0.6),
                 bullets)

    add_slide_number(s, n, total)


def slide_messaging_layer(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, C.MAL["kicker"], C.MAL["title"])

    # Flow diagram (single horizontal pill chain)
    flow_y = Inches(3.15)
    flow_text = C.MAL["flow"]
    add_round_rect(s, T.MARGIN_X, flow_y, Inches(12.13), Inches(0.7),
                   fill=T.BG_DARK_2, line=T.ACCENT_DEEP, line_w=Pt(0.75))
    add_text(s, T.MARGIN_X + Inches(0.4), flow_y, Inches(11.4), Inches(0.7),
             [{"text": flow_text, "size": T.SZ_BODY_S, "color": T.INK_INV,
               "font": T.FONT_MONO}], anchor=MSO_ANCHOR.MIDDLE)

    # Channel rows
    ch_y = Inches(4.1)
    row_h = Inches(0.42)
    row_gap = Inches(0.06)
    STATUS_COLOR = {
        "live": T.GREEN, "queued": T.AMBER, "blocked": T.RED,
        "off-path": T.INK_3, "Sprint 22": T.PURPLE,
    }
    for i, (ch, status, detail) in enumerate(C.MAL["channels"]):
        y = ch_y + i * (row_h + row_gap)
        add_round_rect(s, T.MARGIN_X, y, Inches(12.13), row_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        # Channel name
        add_text(s, T.MARGIN_X + Inches(0.25), y, Inches(3.0), row_h,
                 [{"text": ch, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}],
                 anchor=MSO_ANCHOR.MIDDLE)
        # Status pill
        col = STATUS_COLOR.get(status, T.INK_3)
        add_round_rect(s, T.MARGIN_X + Inches(3.2), y + Inches(0.08),
                       Inches(1.2), row_h - Inches(0.16), fill=col, line=None)
        add_text(s, T.MARGIN_X + Inches(3.2), y + Inches(0.08),
                 Inches(1.2), row_h - Inches(0.16),
                 [{"text": status.upper(), "size": T.SZ_TINY, "color": T.INK_INV,
                   "bold": True, "tracking": 160, "font": T.FONT_TEXT}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Detail
        add_text(s, T.MARGIN_X + Inches(4.6), y, Inches(7.5), row_h,
                 [{"text": detail, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}], anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


def slide_channel_coverage(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Channel coverage", "Telegram-led today. Omnichannel by Phase 22.")

    # 4-col grid
    grid_y = Inches(3.1)
    cols = 4
    cell_w = Inches(3.0)
    cell_h = Inches(1.65)
    gap_x = Inches(0.13)
    gap_y = Inches(0.15)
    for idx, (ch, status, descr, live) in enumerate(C.CHANNEL_COVERAGE):
        r, c = divmod(idx, cols)
        x = T.MARGIN_X + c * (cell_w + gap_x)
        y = grid_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, x, y, cell_w, Inches(0.07),
                 fill=T.GREEN if live else T.AMBER, line=None)
        add_text(s, x + Inches(0.18), y + Inches(0.2), cell_w - Inches(0.3), Inches(0.34),
                 [{"text": ch, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.18), y + Inches(0.55), cell_w - Inches(0.3), Inches(0.3),
                 [{"text": status, "size": T.SZ_TINY,
                   "color": T.GREEN if live else T.AMBER,
                   "bold": True, "tracking": 200, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.18), y + Inches(0.85), cell_w - Inches(0.3), cell_h - Inches(0.95),
                 [{"text": descr, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_market(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Market", "A$2.9B AU SAM. US$30B+ global. Bottom-up to A$120M SOM.")

    # 3 size tiles
    tiles_y = Inches(3.1)
    cell_w = Inches(3.85)
    cell_h = Inches(2.1)
    gap_x = Inches(0.18)
    for i, (lbl, mkt, val, det) in enumerate(C.MARKET_SIZE):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, tiles_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.3), tiles_y + Inches(0.18), cell_w - Inches(0.4), Inches(0.3),
                 [{"text": lbl.upper(), "size": T.SZ_TINY, "color": T.INK_3,
                   "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.3), tiles_y + Inches(0.5), cell_w - Inches(0.4), Inches(0.4),
                 [{"text": mkt, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.3), tiles_y + Inches(0.85), cell_w - Inches(0.4), Inches(0.7),
                 [{"text": val, "size": T.SZ_METRIC_S, "color": T.ACCENT, "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.3), tiles_y + Inches(1.4), cell_w - Inches(0.4), cell_h - Inches(1.5),
                 [{"text": det, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Build blocks
    bb_y = Inches(5.4)
    add_text(s, T.MARGIN_X, bb_y, Inches(8), Inches(0.3),
             [{"text": "BOTTOM-UP BUILD", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for kicker, line in C.MARKET_BUILD_BLOCKS:
        bullets.append([
            {"text": kicker + "  ", "size": T.SZ_BODY_S, "color": T.ACCENT, "bold": True},
            {"text": line, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(6)},
        ])
    add_text(s, T.MARGIN_X, bb_y + Inches(0.3), Inches(12.13), Inches(1.5), bullets)

    add_slide_number(s, n, total)


def slide_competitive(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Competitive landscape", "Vertical depth × AI-native. We sit alone in the top-right.")

    # 2D quadrant chart (X = vertical depth, Y = AI-native, top-right is good)
    plot_x = T.MARGIN_X + Inches(0.5)
    plot_y = Inches(3.1)
    plot_w = Inches(7.5)
    plot_h = Inches(3.7)

    # Plot frame
    add_rect(s, plot_x, plot_y, plot_w, plot_h, fill=T.BG_PANEL, line=T.DIVIDER)

    # Mid grid lines
    add_rect(s, plot_x + plot_w / 2, plot_y, Emu(int(plot_w / 1) * 0 + 9525), plot_h,
             fill=T.DIVIDER, line=None)  # vertical mid
    add_rect(s, plot_x, plot_y + plot_h / 2, plot_w, Pt(0.75), fill=T.DIVIDER, line=None)

    # Axis labels
    add_text(s, plot_x, plot_y + plot_h + Inches(0.06), plot_w, Inches(0.3),
             [{"text": "VERTICAL DEPTH  →", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 200, "font": T.FONT_TEXT}],
             align=PP_ALIGN.CENTER)
    add_text(s, plot_x - Inches(0.5), plot_y, Inches(0.5), plot_h,
             [{"text": "AI-NATIVE  →", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 200, "font": T.FONT_TEXT}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Plot competitors
    for name, vert, ai, tag in C.COMPETITORS:
        # 0..100 → position inside plot
        cx = plot_x + Emu(int(vert / 100 * plot_w))
        cy = plot_y + plot_h - Emu(int(ai / 100 * plot_h))
        is_us = name == "BookedAI"
        size = Inches(0.5) if is_us else Inches(0.32)
        color = T.ACCENT if is_us else T.INK_3
        # Centered dot
        add_round_rect(s, cx - size / 2, cy - size / 2, size, size,
                       fill=color, line=None)
        # Label
        add_text(s, cx + Inches(0.32), cy - Inches(0.18), Inches(2.2), Inches(0.4),
                 [
                     {"text": name, "size": T.SZ_CAPTION,
                      "color": T.ACCENT if is_us else T.INK, "bold": True},
                 ])

    # Right-side legend / commentary
    legend_x = plot_x + plot_w + Inches(0.4)
    legend_w = Inches(4.1)
    add_text(s, legend_x, plot_y, legend_w, Inches(0.32),
             [{"text": "WHO COMPETES — AND WHERE", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for name, _, _, tag in C.COMPETITORS:
        bullets.append([
            {"text": f"{name}: ", "size": T.SZ_BODY_S,
             "color": T.ACCENT if name == "BookedAI" else T.INK, "bold": True},
            {"text": tag, "size": T.SZ_BODY_S, "color": T.INK_2,
             "font": T.FONT_TEXT, "space_after": Pt(8)},
        ])
    add_text(s, legend_x, plot_y + Inches(0.4), legend_w, plot_h - Inches(0.4), bullets)

    add_slide_number(s, n, total)


def slide_defensibility(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Defensibility", "Three moats compound as we scale.")

    # Lede
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(12.13), Inches(0.7),
             [{"text": C.DEFENSIBILITY_LEDE, "size": T.SZ_BODY, "color": T.INK_2,
               "font": T.FONT_TEXT}])

    # Three moat columns
    moat_y = Inches(3.95)
    cell_w = Inches(3.85)
    cell_h = Inches(3.0)
    gap_x = Inches(0.18)
    for i, m in enumerate(C.DEFENSIBILITY):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, moat_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, x, moat_y, cell_w, Inches(0.08), fill=T.ACCENT, line=None)
        add_text(s, x + Inches(0.3), moat_y + Inches(0.25), cell_w - Inches(0.4), Inches(0.3),
                 [{"text": m["kicker"].upper(), "size": T.SZ_TINY, "color": T.ACCENT,
                   "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.3), moat_y + Inches(0.55), cell_w - Inches(0.4), Inches(0.7),
                 [{"text": m["title"], "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.3), moat_y + Inches(1.25), cell_w - Inches(0.4), cell_h - Inches(1.4),
                 [{"text": m["body"], "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_business_model(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Business Model", C.BUSINESS_MODEL["headline"])

    # Lede
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(12.13), Inches(0.6),
             [{"text": C.BUSINESS_MODEL["lede"], "size": T.SZ_BODY, "color": T.INK_2,
               "font": T.FONT_TEXT}])

    # Three components
    comp_y = Inches(3.7)
    cell_w = Inches(3.85)
    cell_h = Inches(1.8)
    gap_x = Inches(0.18)
    for i, (name, freq, desc) in enumerate(C.BUSINESS_MODEL["components"]):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, comp_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.28), comp_y + Inches(0.18), cell_w - Inches(0.4), Inches(0.36),
                 [{"text": name, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.28), comp_y + Inches(0.52), cell_w - Inches(0.4), Inches(0.3),
                 [{"text": freq.upper(), "size": T.SZ_TINY, "color": T.ACCENT, "bold": True,
                   "tracking": 200, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.28), comp_y + Inches(0.85), cell_w - Inches(0.4), cell_h - Inches(0.95),
                 [{"text": desc, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Tier table
    tier_y = Inches(5.7)
    add_text(s, T.MARGIN_X, tier_y, Inches(8), Inches(0.3),
             [{"text": "PERSONA TIERS", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    tier_y2 = tier_y + Inches(0.3)
    row_h = Inches(0.36)
    for i, (persona, price, descr) in enumerate(C.BUSINESS_MODEL["tiers"]):
        y = tier_y2 + i * row_h
        add_text(s, T.MARGIN_X, y, Inches(2.6), row_h,
                 [{"text": persona, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(2.6), y, Inches(2.0), row_h,
                 [{"text": price, "size": T.SZ_BODY_S, "color": T.ACCENT, "bold": True,
                   "font": T.FONT_MONO}], anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(4.6), y, Inches(7.5), row_h,
                 [{"text": descr, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


def slide_unit_economics(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Unit Economics", "Setup fee covers CAC. Commission compounds the LTV.")

    # 4 metric tiles
    tiles_y = Inches(3.1)
    cell_w = Inches(2.9)
    cell_h = Inches(2.0)
    gap_x = Inches(0.16)
    for i, (val, label, sub) in enumerate(C.UNIT_ECONOMICS):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, tiles_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.25), tiles_y + Inches(0.2), cell_w - Inches(0.4), Inches(0.7),
                 [{"text": val, "size": T.SZ_METRIC_S, "color": T.ACCENT,
                   "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.25), tiles_y + Inches(0.85), cell_w - Inches(0.4), Inches(0.34),
                 [{"text": label, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.25), tiles_y + Inches(1.2), cell_w - Inches(0.4), cell_h - Inches(1.3),
                 [{"text": sub, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    # Note
    note_y = Inches(5.4)
    add_round_rect(s, T.MARGIN_X, note_y, Inches(12.13), Inches(1.4),
                   fill=T.BG_DARK_2, line=None)
    add_text(s, T.MARGIN_X + Inches(0.4), note_y + Inches(0.18), Inches(11.4), Inches(0.32),
             [{"text": "WHY THIS WORKS", "size": T.SZ_TINY, "color": T.ACCENT_SOFT,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    add_text(s, T.MARGIN_X + Inches(0.4), note_y + Inches(0.52), Inches(11.4), Inches(0.85),
             [{"text": C.UNIT_ECONOMICS_NOTE, "size": T.SZ_BODY_S, "color": T.INK_INV,
               "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_traction(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Traction", C.TRACTION["headline"])

    # Two columns: Shipped, Next
    col_w = Inches(5.95)
    col_h = Inches(4.0)
    col_y = Inches(2.95)
    gap = Inches(0.23)

    # Shipped
    add_round_rect(s, T.MARGIN_X, col_y, col_w, col_h, fill=T.BG_PANEL, line=T.DIVIDER)
    add_rect(s, T.MARGIN_X, col_y, col_w, Inches(0.07), fill=T.GREEN, line=None)
    add_text(s, T.MARGIN_X + Inches(0.32), col_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
             [{"text": "SHIPPED — production today",
               "size": T.SZ_TINY, "color": T.GREEN, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    bullets = []
    for k, v in C.TRACTION["shipped"]:
        bullets.append([
            {"text": "✓  ", "size": T.SZ_BODY_S, "color": T.GREEN, "bold": True},
            {"text": k + " — ", "size": T.SZ_BODY_S, "color": T.INK, "bold": True},
            {"text": v, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(8)},
        ])
    add_text(s, T.MARGIN_X + Inches(0.32), col_y + Inches(0.55),
             col_w - Inches(0.5), col_h - Inches(0.7), bullets)

    # Next
    nx = T.MARGIN_X + col_w + gap
    add_round_rect(s, nx, col_y, col_w, col_h, fill=T.BG_PANEL, line=T.DIVIDER)
    add_rect(s, nx, col_y, col_w, Inches(0.07), fill=T.ACCENT, line=None)
    add_text(s, nx + Inches(0.32), col_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
             [{"text": "NEXT 60 DAYS — Sprints 19-22",
               "size": T.SZ_TINY, "color": T.ACCENT, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    bullets = []
    for k, v in C.TRACTION["next"]:
        bullets.append([
            {"text": "›  ", "size": T.SZ_BODY_S, "color": T.ACCENT, "bold": True},
            {"text": k + " — ", "size": T.SZ_BODY_S, "color": T.INK, "bold": True},
            {"text": v, "size": T.SZ_BODY_S, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(8)},
        ])
    add_text(s, nx + Inches(0.32), col_y + Inches(0.55),
             col_w - Inches(0.5), col_h - Inches(0.7), bullets)

    add_slide_number(s, n, total)


def slide_roadmap(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Roadmap", "Phase 17 → Phase 22+ · revenue milestones, not feature lists.")

    # Vertical list of phases
    rows_y = Inches(3.05)
    row_h = Inches(0.78)
    row_gap = Inches(0.06)
    for i, (phase, dt, outcome, mid, mtitle, line) in enumerate(C.ROADMAP_PHASES):
        y = rows_y + i * (row_h + row_gap)
        add_round_rect(s, T.MARGIN_X, y, Inches(12.13), row_h,
                       fill=T.BG_PANEL, line=T.DIVIDER)
        add_rect(s, T.MARGIN_X, y, Inches(0.08), row_h, fill=T.ACCENT, line=None)
        # Phase + date
        add_text(s, T.MARGIN_X + Inches(0.3), y + Inches(0.1), Inches(2.5), Inches(0.3),
                 [{"text": phase, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}])
        add_text(s, T.MARGIN_X + Inches(0.3), y + Inches(0.4), Inches(2.5), Inches(0.3),
                 [{"text": dt, "size": T.SZ_CAPTION, "color": T.ACCENT,
                   "font": T.FONT_MONO}])
        # Outcome pill
        add_round_rect(s, T.MARGIN_X + Inches(2.85), y + Inches(0.18),
                       Inches(2.4), Inches(0.42), fill=T.ACCENT_DEEP, line=None)
        add_text(s, T.MARGIN_X + Inches(2.85), y + Inches(0.18),
                 Inches(2.4), Inches(0.42),
                 [{"text": outcome, "size": T.SZ_CAPTION, "color": T.INK_INV, "bold": True}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Milestone title
        add_text(s, T.MARGIN_X + Inches(5.4), y + Inches(0.08), Inches(6.7), Inches(0.34),
                 [
                     {"text": f"{mid}  ", "size": T.SZ_TINY, "color": T.INK_3, "bold": True,
                      "tracking": 200, "font": T.FONT_MONO},
                     {"text": mtitle, "size": T.SZ_BODY_S, "color": T.INK, "bold": True},
                 ])
        # Revenue line
        add_text(s, T.MARGIN_X + Inches(5.4), y + Inches(0.4), Inches(6.7), Inches(0.36),
                 [{"text": line, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_team(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Team", "Founder-led today. Hiring plan locked through Sprint 22.")

    # Team cards (2)
    team_y = Inches(2.95)
    cell_w = Inches(5.95)
    cell_h = Inches(1.8)
    gap_x = Inches(0.23)
    for i, m in enumerate(C.TEAM):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, team_y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        # Avatar circle
        add_round_rect(s, x + Inches(0.25), team_y + Inches(0.25),
                       Inches(0.85), Inches(0.85),
                       fill=T.ACCENT_DEEP, line=None)
        # Initials
        initials = "".join(p[0] for p in m["name"].split() if p)[:2].upper()
        add_text(s, x + Inches(0.25), team_y + Inches(0.25),
                 Inches(0.85), Inches(0.85),
                 [{"text": initials, "size": T.SZ_BODY, "color": T.INK_INV, "bold": True}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name + role + bio
        add_text(s, x + Inches(1.3), team_y + Inches(0.2), cell_w - Inches(1.5), Inches(0.4),
                 [{"text": m["name"], "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(1.3), team_y + Inches(0.6), cell_w - Inches(1.5), Inches(0.32),
                 [{"text": m["role"], "size": T.SZ_TINY, "color": T.ACCENT,
                   "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(1.3), team_y + Inches(0.95), cell_w - Inches(1.5), cell_h - Inches(1.05),
                 [{"text": m["bio"], "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}])

    # Hiring plan
    hire_y = Inches(5.0)
    add_text(s, T.MARGIN_X, hire_y, Inches(8), Inches(0.3),
             [{"text": "HIRING PLAN — funded by this round",
               "size": T.SZ_TINY, "color": T.INK_3, "bold": True, "tracking": 220,
               "font": T.FONT_TEXT}])
    rows_y = hire_y + Inches(0.32)
    row_h = Inches(0.42)
    for i, (role, dt, scope) in enumerate(C.HIRING_PLAN):
        y = rows_y + i * row_h
        add_text(s, T.MARGIN_X, y, Inches(3.0), row_h,
                 [{"text": role, "size": T.SZ_BODY_S, "color": T.INK, "bold": True}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(3.0), y, Inches(2.6), row_h,
                 [{"text": dt, "size": T.SZ_CAPTION, "color": T.ACCENT,
                   "bold": True, "font": T.FONT_MONO}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, T.MARGIN_X + Inches(5.6), y, Inches(6.5), row_h,
                 [{"text": scope, "size": T.SZ_CAPTION, "color": T.INK_2,
                   "font": T.FONT_TEXT}], anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, n, total)


def slide_partners(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Partners & Integrations", "Production stack — vendors integrated, not promised.")

    # 4-col grid of partner tiles
    grid_y = Inches(3.0)
    cols = 4
    cell_w = Inches(3.0)
    cell_h = Inches(1.05)
    gap_x = Inches(0.13)
    gap_y = Inches(0.13)
    for idx, (name, role) in enumerate(C.PARTNERS[: cols * 3]):
        r, c = divmod(idx, cols)
        x = T.MARGIN_X + c * (cell_w + gap_x)
        y = grid_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=T.BG_PANEL, line=T.DIVIDER)
        add_text(s, x + Inches(0.2), y + Inches(0.16), cell_w - Inches(0.3), Inches(0.36),
                 [{"text": name, "size": T.SZ_BODY, "color": T.INK, "bold": True}])
        add_text(s, x + Inches(0.2), y + Inches(0.56), cell_w - Inches(0.3), cell_h - Inches(0.6),
                 [{"text": role, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT}])

    add_slide_number(s, n, total)


def slide_financial(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "Financial projection", "Bottom-up. Tied to phase milestones, not vibes.")

    # Header row
    header_y = Inches(3.05)
    headers = ["Year", "Tenants", "Blended ARPU", "ARR", "Gross margin", "Team posture"]
    col_widths = [Inches(2.0), Inches(1.4), Inches(2.0), Inches(1.7), Inches(1.7), Inches(3.33)]
    cur_x = T.MARGIN_X
    add_round_rect(s, T.MARGIN_X, header_y, Inches(12.13), Inches(0.42),
                   fill=T.BG_DARK_2, line=None)
    for i, h in enumerate(headers):
        add_text(s, cur_x + Inches(0.18), header_y, col_widths[i], Inches(0.42),
                 [{"text": h.upper(), "size": T.SZ_TINY, "color": T.INK_INV_2,
                   "bold": True, "tracking": 200, "font": T.FONT_TEXT}],
                 anchor=MSO_ANCHOR.MIDDLE)
        cur_x = cur_x + col_widths[i]

    # Body rows
    row_y = header_y + Inches(0.5)
    row_h = Inches(0.55)
    for i, row in enumerate(C.FINANCIAL_PROJECTION["rows"]):
        y = row_y + i * (row_h + Inches(0.05))
        bg = T.BG_PANEL if i % 2 == 0 else T.BG_LIGHT
        add_round_rect(s, T.MARGIN_X, y, Inches(12.13), row_h, fill=bg, line=T.DIVIDER, line_w=Pt(0.5))
        cur_x = T.MARGIN_X
        for j, cell in enumerate(row):
            color = T.ACCENT if j == 3 else T.INK
            bold = (j in (0, 3))
            add_text(s, cur_x + Inches(0.18), y, col_widths[j] - Inches(0.18), row_h,
                     [{"text": cell, "size": T.SZ_BODY_S, "color": color, "bold": bold,
                       "font": T.FONT_DISPLAY if j != 3 else T.FONT_MONO}],
                     anchor=MSO_ANCHOR.MIDDLE)
            cur_x = cur_x + col_widths[j]

    # Assumptions
    assum_y = Inches(5.7)
    add_text(s, T.MARGIN_X, assum_y, Inches(8), Inches(0.3),
             [{"text": "ASSUMPTIONS", "size": T.SZ_TINY, "color": T.INK_3,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for line in C.FINANCIAL_PROJECTION["assumptions"]:
        bullets.append([
            {"text": "•  ", "size": T.SZ_CAPTION, "color": T.ACCENT, "bold": True},
            {"text": line, "size": T.SZ_CAPTION, "color": T.INK_2, "font": T.FONT_TEXT,
             "space_after": Pt(4)},
        ])
    add_text(s, T.MARGIN_X, assum_y + Inches(0.3), Inches(12.13), Inches(1.4), bullets)

    add_slide_number(s, n, total)


def slide_ask(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(s, "The Ask", C.ASK["headline"], dark=True)

    # Structure line
    add_text(s, T.MARGIN_X, Inches(2.95), Inches(12.13), Inches(0.5),
             [{"text": C.ASK["structure"], "size": T.SZ_SUBTITLE, "color": T.ACCENT_SOFT,
               "font": T.FONT_TEXT}])

    # Use of funds — 4 columns
    uof_y = Inches(3.7)
    cell_w = Inches(2.9)
    cell_h = Inches(1.8)
    gap_x = Inches(0.16)
    for i, (pct, label, descr) in enumerate(C.ASK["use_of_funds"]):
        x = T.MARGIN_X + i * (cell_w + gap_x)
        add_round_rect(s, x, uof_y, cell_w, cell_h, fill=T.BG_PANEL_DARK,
                       line=T.DIVIDER_DARK)
        add_text(s, x + Inches(0.25), uof_y + Inches(0.18), cell_w - Inches(0.4), Inches(0.7),
                 [{"text": pct, "size": T.SZ_METRIC_S, "color": T.ACCENT_SOFT,
                   "bold": True, "tracking": -20}])
        add_text(s, x + Inches(0.25), uof_y + Inches(0.85), cell_w - Inches(0.4), Inches(0.34),
                 [{"text": label, "size": T.SZ_BODY_S, "color": T.INK_INV, "bold": True}])
        add_text(s, x + Inches(0.25), uof_y + Inches(1.2), cell_w - Inches(0.4), cell_h - Inches(1.3),
                 [{"text": descr, "size": T.SZ_TINY, "color": T.INK_INV_2,
                   "font": T.FONT_TEXT}])

    # Milestones row
    ms_y = Inches(5.7)
    add_text(s, T.MARGIN_X, ms_y, Inches(8), Inches(0.3),
             [{"text": "MILESTONES THIS ROUND FUNDS",
               "size": T.SZ_TINY, "color": T.ACCENT_SOFT,
               "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
    bullets = []
    for line in C.ASK["milestones"]:
        bullets.append([
            {"text": "→  ", "size": T.SZ_BODY_S, "color": T.ACCENT_SOFT, "bold": True},
            {"text": line, "size": T.SZ_BODY_S, "color": T.INK_INV, "font": T.FONT_TEXT,
             "space_after": Pt(4)},
        ])
    add_text(s, T.MARGIN_X, ms_y + Inches(0.3), Inches(12.13), Inches(1.4), bullets)

    add_slide_number(s, n, total, dark=True)


def slide_cta(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s, T.BG_DARK)
    add_brand_mark(s, dark=True)

    # Hero text
    add_text(s, T.MARGIN_X, Inches(2.0), Inches(12.13), Inches(1.6),
             [{"text": C.CTA["headline"], "size": Pt(48), "color": T.INK_INV,
               "bold": True, "tracking": -30}])
    add_text(s, T.MARGIN_X, Inches(3.5), Inches(12.13), Inches(0.6),
             [{"text": C.CTA["subhead"], "size": T.SZ_SUBTITLE, "color": T.INK_INV_2,
               "font": T.FONT_TEXT}])

    # Hairline
    add_rect(s, T.MARGIN_X, Inches(4.3), Inches(1.5), Pt(2), fill=T.ACCENT, line=None)

    # CTA grid
    cta_y = Inches(4.6)
    cell_w = Inches(5.95)
    cell_h = Inches(0.85)
    gap_x = Inches(0.23)
    gap_y = Inches(0.16)
    for idx, (label, value) in enumerate(C.CTA["ctas"]):
        r, c = divmod(idx, 2)
        x = T.MARGIN_X + c * (cell_w + gap_x)
        y = cta_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=T.BG_PANEL_DARK, line=T.DIVIDER_DARK)
        add_text(s, x + Inches(0.3), y + Inches(0.1), cell_w - Inches(0.4), Inches(0.32),
                 [{"text": label, "size": T.SZ_TINY, "color": T.ACCENT_SOFT,
                   "bold": True, "tracking": 220, "font": T.FONT_TEXT}])
        add_text(s, x + Inches(0.3), y + Inches(0.4), cell_w - Inches(0.4), Inches(0.45),
                 [{"text": value, "size": T.SZ_BODY, "color": T.INK_INV, "bold": True,
                   "font": T.FONT_MONO}])

    # Closer
    add_text(s, T.MARGIN_X, Inches(6.7), Inches(12.13), Inches(0.4),
             [{"text": C.CTA["closer"], "size": T.SZ_BODY_S, "color": T.ACCENT_SOFT,
               "bold": True, "font": T.FONT_MONO}], align=PP_ALIGN.CENTER)

    add_slide_number(s, n, total, dark=True)


# ---------- Orchestration ----------

# Slide registry: list of (name, builder) — used to compute total + render.
SLIDES = [
    ("cover", slide_cover),
    ("vision", slide_vision),
    ("problem", slide_problem),
    ("why_now", slide_why_now),
    ("solution", slide_solution),
    ("surfaces", slide_surfaces),
    ("journey", slide_journey),
    ("agents", slide_agents),
    ("live_evidence", slide_live_evidence),
    ("tenant_chess", lambda prs, n, t: slide_tenant_case(
        prs, n, t, C.TENANT_CASES[0], T.TENANT_CHESS)),
    ("tenant_aimentor", lambda prs, n, t: slide_tenant_case(
        prs, n, t, C.TENANT_CASES[1], T.TENANT_AIMENTOR)),
    ("tenant_futureswim", lambda prs, n, t: slide_tenant_case(
        prs, n, t, C.TENANT_CASES[2], T.TENANT_FUTURESWIM)),
    ("architecture", slide_architecture),
    ("tech_stack", slide_tech_stack),
    ("messaging_layer", slide_messaging_layer),
    ("channel_coverage", slide_channel_coverage),
    ("market", slide_market),
    ("competitive", slide_competitive),
    ("defensibility", slide_defensibility),
    ("business_model", slide_business_model),
    ("unit_economics", slide_unit_economics),
    ("traction", slide_traction),
    ("roadmap", slide_roadmap),
    ("team", slide_team),
    ("partners", slide_partners),
    ("financial", slide_financial),
    ("ask", slide_ask),
    ("cta", slide_cta),
]


def build():
    prs = Presentation()
    prs.slide_width = T.SLIDE_WIDTH
    prs.slide_height = T.SLIDE_HEIGHT

    total = len(SLIDES)
    for i, (name, fn) in enumerate(SLIDES, start=1):
        fn(prs, i, total)

    OUTPUT.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT / f"BookedAI_Investor_Pitch_{C.DECK_DATE.isoformat()}.pptx"
    prs.save(str(out_path))

    size_kb = out_path.stat().st_size / 1024
    print(f"✓ Built {out_path.name}  ·  {total} slides  ·  {size_kb:,.0f} KB")
    print(f"  → {out_path}")
    return out_path


if __name__ == "__main__":
    build()
